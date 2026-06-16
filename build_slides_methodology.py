# -*- coding: utf-8 -*-
"""Insert engineering-methodology (formula) slides into the GeoTrack defence deck.

Formulas are taken directly from the implemented source code:
  - mobile SetBoundaryBottomSheet.tsx -> N-sample GPS averaging
  - GeofenceService.php   -> Haversine distance, circular & polygon (ray-cast) geofence
  - FaceMatchService.php  -> grayscale luminance, average hash, Hamming distance,
                             confidence, AWS Rekognition cosine similarity
  - AttendanceController.php  -> lateness, present/late classification, composite gate
  - PresenceCheckController.php -> randomised presence checks, present-throughout

Each new slide is a CLONE of the existing "Methodology" slide (slide 13) so the
running header, title style, fonts and theme colours match the rest of the deck.
Only the title and body text are rewritten.
"""
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

SRC = "docs/Essien Abasiama (20201742) Post-Data Slide.pptx"
TITLE_TEXT = "Methodology"               # text that identifies the title shape
HEADER_TEXT = "Materials and Methods"    # running header label (kept as-is)
GRAY = RGBColor(0x66, 0x66, 0x66)
TNR = "Times New Roman"

prs = Presentation(SRC)
METH_INDEX = 12  # zero-based -> slide 13 "Methodology"


# ── slide cloning + ordering ────────────────────────────────────────────────
def clone_slide(prs, index):
    """Deep-copy every shape of slide `index` onto a new BLANK slide."""
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    # strip placeholders the layout auto-inserts
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    return dest


def move_slide(prs, old_index, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[old_index]
    lst.remove(el)
    lst.insert(new_index, el)


# ── text helpers ────────────────────────────────────────────────────────────
def find_shapes(slide):
    """Return (title_shape, body_shape) by inspecting current text."""
    title = body = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if txt == TITLE_TEXT and title is None:
            title = sh
        elif txt.startswith("Adopted a Design"):
            body = sh
    return title, body


def set_title(shape, text):
    """Replace title text and resize so the longer methodology titles fit on
    one line (the original box only held the single word 'Methodology')."""
    # widen the box and pull it left so it spans most of the slide width
    shape.left = Pt(0.5 * 72)     # 0.5"
    shape.width = Pt(9.0 * 72)    # 9.0"
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    runs = p.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
        runs[0].font.size = Pt(22)   # down from 36pt so it stays one line
    else:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(22)


def build_body(shape, lines):
    """Rebuild a body text frame from a list of line specs.

    Each spec: dict(text, size, bold=False, italic=False,
                    align='left'|'center', color=None, eqno=None, space_before=0)
    `eqno`, when present, is appended in gray parentheses at the line end.
    """
    tf = shape.text_frame
    tf.word_wrap = True
    # clear all existing paragraphs
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    for i, spec in enumerate(lines):
        p = tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
        }[spec.get("align", "left")]
        p.space_after = Pt(spec.get("space_after", 4))
        p.space_before = Pt(spec.get("space_before", 0))

        r = p.add_run()
        r.text = spec["text"]
        f = r.font
        f.name = TNR
        f.size = Pt(spec["size"])
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        if spec.get("color"):
            f.color.rgb = spec["color"]

        if spec.get("eqno"):
            rn = p.add_run()
            rn.text = "   (" + spec["eqno"] + ")"
            rn.font.name = TNR
            rn.font.size = Pt(spec["size"] - 2)
            rn.font.italic = True
            rn.font.color.rgb = GRAY


def L(text, **kw):
    kw["text"] = text
    kw.setdefault("size", 15)
    return kw


# ── slide content ───────────────────────────────────────────────────────────
SLIDES = [
    # ---- 1. GPS positioning & Haversine -------------------------------------
    {
        "title": "Methodology: GPS Positioning & Geofencing (I)",
        "body": [
            L("A student's position P = (φₚ, λₚ) is sampled N = 5 times at the "
              "highest GNSS accuracy and averaged to suppress random satellite noise:",
              size=15, space_after=6),
            L("P̄ = ( (1/N)·Σ φᵢ ,  (1/N)·Σ λᵢ )",
              size=18, align="center", bold=True, eqno="3.1", space_after=8),
            L("The great-circle distance to the venue centre uses the Haversine formula:",
              size=15, space_after=6),
            L("a = sin²(Δφ/2) + cos φ₁ · cos φ₂ · sin²(Δλ/2)",
              size=17, align="center", eqno="3.2", space_after=4),
            L("c = 2 · atan2( √a , √(1 − a) )",
              size=17, align="center", eqno="3.3", space_after=4),
            L("d = R · c ,    R = 6 371 000 m  (mean Earth radius)",
              size=17, align="center", bold=True, eqno="3.4", space_after=6),
            L("where d is the distance (m) between the student and the venue centre; "
              "φ, λ are latitude/longitude in radians.",
              size=13, italic=True, color=GRAY, space_after=0),
        ],
    },
    # ---- 2. Boundary decision -----------------------------------------------
    {
        "title": "Methodology: Geofencing — Boundary Decision (II)",
        "body": [
            L("Circular venue (centre C = (φ₀, λ₀), radius r): the student is admitted "
              "if and only if the Haversine distance is within the radius:",
              size=15, space_after=6),
            L("Inside  ⇔  d ≤ r",
              size=18, align="center", bold=True, eqno="3.5", space_after=8),
            L("Polygon venue: the ray-casting (even–odd) rule. A horizontal ray is cast "
              "from P; for each edge (Vᵢ, Vⱼ) the crossing test is:",
              size=15, space_after=6),
            L("(yᵢ > yₚ) ≠ (yⱼ > yₚ)  ∧  xₚ < (xⱼ − xᵢ)(yₚ − yᵢ)/(yⱼ − yᵢ) + xᵢ",
              size=15, align="center", eqno="3.6", space_after=8),
            L("An ODD number of edge crossings ⇒ the point lies inside the boundary.",
              size=15, bold=True, space_after=6),
            L("where (x, y) = (longitude, latitude); the reported GNSS accuracy is stored "
              "with each record so near-boundary decisions can be audited.",
              size=13, italic=True, color=GRAY, space_after=0),
        ],
    },
    # ---- 3. Face verification -----------------------------------------------
    {
        "title": "Methodology: Biometric (Face) Verification",
        "body": [
            L("A live selfie is reduced to a descriptor and compared with the student's "
              "enrolled reference face.", size=14, space_after=4),
            L("(a)  On-device perceptual-hash matcher:", size=14, bold=True, space_after=2),
            L("Y = 0.299R + 0.587G + 0.114B        (grayscale luminance)",
              size=15, align="center", eqno="3.7", space_after=3),
            L("μ = (1/N²)·Σ Yᵢ ;    bᵢ = 1 if Yᵢ ≥ μ, else 0     (256-bit hash)",
              size=15, align="center", eqno="3.8", space_after=3),
            L("D(Hᵣ, Hₗ) = Σ (bᵣ,ᵢ ⊕ bₗ,ᵢ) ;   Confidence = 100·(1 − D/256) %",
              size=15, align="center", eqno="3.9", space_after=6),
            L("(b)  Deep-CNN matcher (AWS Rekognition embeddings f ∈ ℝⁿ):",
              size=14, bold=True, space_after=2),
            L("S(fᵣ, fₗ) = (fᵣ · fₗ) / (‖fᵣ‖ ‖fₗ‖)        (cosine similarity)",
              size=15, align="center", eqno="3.10", space_after=3),
            L("Verified  ⇔  S × 100 ≥ T        (acceptance threshold T = 70%)",
              size=15, align="center", bold=True, eqno="3.11", space_after=0),
        ],
    },
    # ---- 4. Binding, timing & composite gate --------------------------------
    {
        "title": "Methodology: Binding, Timing & Composite Decision",
        "body": [
            L("Device binding (one account ↔ one phone) — every sensitive request must "
              "carry the bound device fingerprint:", size=14, space_after=4),
            L("Allow  ⇔  Uᵣₑǫ = Uᵇₒᵤₙᵈ  ∧  revoked_at = NULL",
              size=16, align="center", bold=True, eqno="3.12", space_after=6),
            L("Lateness and status inside the active window [tₛ, tₑ]:", size=14, space_after=3),
            L("m = t − tₛ ;    Status = PRESENT if m ≤ τ, else LATE",
              size=16, align="center", eqno="3.13", space_after=6),
            L("Continuous presence (randomised checks every interval I); a record is "
              "valid throughout when no check is missed:", size=14, space_after=3),
            L("PresentThroughout  ⇔  M = 0",
              size=16, align="center", eqno="3.14", space_after=6),
            L("Final decision — all gates must hold simultaneously:", size=14, bold=True, space_after=3),
            L("Granted = B ∧ G ∧ F ∧ W",
              size=18, align="center", bold=True, eqno="3.15", space_after=4),
            L("B = bound device, G = inside geofence, F = face verified, W = within time window.",
              size=12, italic=True, color=GRAY, space_after=0),
        ],
    },
]


# ── build ───────────────────────────────────────────────────────────────────
insert_at = METH_INDEX + 1  # right after the existing Methodology slide
for offset, spec in enumerate(SLIDES):
    slide = clone_slide(prs, METH_INDEX)
    title_sh, body_sh = find_shapes(slide)
    if title_sh is None or body_sh is None:
        raise SystemExit(f"Could not locate title/body shapes on clone {offset}")
    set_title(title_sh, spec["title"])
    build_body(body_sh, spec["body"])
    # new slide is appended at the end; move it into place
    move_slide(prs, len(prs.slides._sldIdLst) - 1, insert_at + offset)

prs.save(SRC)
print(f"Inserted {len(SLIDES)} methodology formula slides after slide {METH_INDEX + 1}.")
print(f"Deck now has {len(prs.slides._sldIdLst)} slides.")
